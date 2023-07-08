from __future__ import annotations

__all__ = ["Patch", "PatchStitching", "MultiPatchRegistration"]

import math
from typing import Iterable, List, Literal, Optional, Tuple

import cv2
import networkx as nx
import numpy as np
import SimpleITK as sitk
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from tqdm import tqdm

from ..utilities.geometry import Point, Rect

BlendMode = Literal["average", "max", "min", "radial", "gaussian", "none"]


class Patch:
    def __init__(self, image: np.ndarray, domain: Rect, alias: str = None):
        self.image = image
        self.domain = domain
        self.alias = alias

    def duplicate(self) -> Patch:
        return Patch(self.image, self.domain, alias=self.alias)

    @staticmethod
    def from_picture(picture: Picture):
        buffer = np.frombuffer(picture.image.blob, dtype=np.uint8)
        img = cv2.imdecode(buffer, flags=1).astype(np.float32) / 255
        domain = Rect(x=picture.left, y=picture.top, h=picture.height, w=picture.width)

        if picture.crop_bottom != 0 or picture.crop_left != 0 or picture.crop_right != 0 or picture.crop_top != 0:
            assert (
                picture.crop_bottom >= 0
                and picture.crop_left >= 0
                and picture.crop_right >= 0
                and picture.crop_top >= 0
            ), (
                "Negative crop values are not supported:"
                f"{picture.crop_left}, {picture.crop_top}, {picture.crop_right}, {picture.crop_bottom}"
            )
            h, w = img.shape[:2]
            x0 = int(picture.crop_left * w)
            x1 = int((1 - picture.crop_right) * w)
            y0 = int(math.ceil(picture.crop_top * h))
            y1 = int(math.ceil((1 - picture.crop_bottom) * h))
            img = img[y0:y1, x0:x1]

        return Patch(img, domain, alias=picture.name)

    def resized_image(self, target_shape: Optional[Tuple[int, int] | Rect] = None):
        if target_shape is None:
            target_shape = self.domain.shape
        elif isinstance(target_shape, Rect):
            target_shape = target_shape.to_int().shape
        if self.image.shape[:2] != target_shape:
            return cv2.resize(self.image, target_shape[::-1])
        return self.image

    @property
    def resolution(self):
        return self.image.shape[0] / self.domain.h

    def extract_domain(self, domain: Rect) -> Patch:
        """
        Extract a sub-domain from the patch.

        Parameters
        ----------
        domain
            The sub-domain to extract.

        Returns
        -------
            The sub-domain of the patch.
        """
        assert not Rect.is_empty(domain), "The domain to extract must be not empty."
        assert domain in self.domain, "The domain to extract must be included in the patch domain."
        image_domain = ((domain - self.domain.top_left) * self.resolution).to_int()
        return Patch(self.image[image_domain.slice()], domain, alias=self.alias)


class PatchStitching(list):
    def __init__(self, iterable: Iterable[Patch] = ()):
        super().__init__((p.duplicate() for p in iterable))
        self.shape = None

    def copy(self) -> PatchStitching:
        return PatchStitching(self)

    def __getitem__(self, item: int | slice[int]) -> Patch | list[Patch]:
        return super().__getitem__(item)

    @staticmethod
    def load_pptx_slide(pptx_path: str, slide_id: int = 0):
        prs = Presentation(pptx_path)
        return PatchStitching.from_pptx_slide(prs.slides[slide_id])

    @staticmethod
    def from_pptx_slide(slide: Slide):
        def find_pictures_recursive(group):
            pictures: list[Patch] = []
            for s in group:
                if isinstance(s, GroupShape):
                    if s.shapes:
                        # Find pictures in group
                        group_pictures = find_pictures_recursive(s.shapes)

                        # Transform the group pictures to take into account its position and size
                        group_domain = Rect.union(patch.domain for patch in group_pictures)
                        group_boundaries = Rect(x=s.left, y=s.top, w=s.width, h=s.height)
                        if group_domain != group_boundaries:
                            domain_transform = group_domain.transforms_to(group_boundaries)
                            for patch in group_pictures:
                                patch.domain = domain_transform(patch.domain)

                        # Add the group pictures to list
                        pictures += group_pictures
                elif isinstance(s, Picture):
                    pictures += [Patch.from_picture(s)]
            return pictures

        pictures = find_pictures_recursive(slide.shapes)
        stitching = PatchStitching(pictures)

        return stitching

    def crop_to_content(self, inplace=True) -> Tuple[Point, List[Rect]]:
        """
        Compute the patches domains so no empty space is left around them when they are stitched together and so their
        resolution is equal to the finest resolution.

        Parameters
        ----------
        inplace
            Whether to modify the patches domains in place or not. (default: True)

        Returns
        -------
            The shape of the stitched image and the domain of each patch.
        """
        resolutions = {patch.resolution for patch in self}
        finest_resolution = max(resolutions)

        domain = self.domain
        patch_domains = []
        for patch in self:
            patch_domain = patch.domain
            patch_domain -= domain.top_left
            patch_domain *= finest_resolution
            patch_domain = patch_domain.to_int()

            patch_domains.append(patch_domain)
            if inplace:
                patch.domain = patch_domain

        shape = tuple(Rect.union(patch_domains).shape)

        if inplace:
            self.shape = shape
        return shape, patch_domains

    @property
    def domain(self):
        return Rect.union(patch.domain for patch in self)

    def extract_domain(self, domain: Rect) -> PatchStitching:
        """
        Extract a sub-domain from the patch stitching.

        Parameters
        ----------
        domain
            The sub-domain to extract.

        Returns
        -------
            A new patch stitching with the sub-domain of each patch.
        """
        assert domain in self.domain, "The domain to extract must be included in the patch stitching domain."
        patches = []
        for patch in self:
            patch_domain = domain & patch.domain
            if patch_domain:
                patches.append(patch.extract_domain(patch_domain))
        return PatchStitching(patches)

    def map_images(self, fn, verbose=False) -> PatchStitching:
        """
        Apply a function to each patch image.

        Parameters
        ----------
        fn
            The function to apply to each patch image.

        Returns
        -------
            A new patch stitching with the modified patches.
        """
        patches = PatchStitching()
        for patch in tqdm(self, desc="Map images", disable=not verbose):
            patch = patch.duplicate()
            patch.image = fn(patch.image)
            patches.append(patch)
        return patches

    def patch_intersection_graph(self, threshold: float = 0):
        """
        Generate the optimal dependency graph between patches based on their intersection area.
        """
        # - Compute the area of intersection between each pair of patches
        nb_patches = len(self)
        areas = np.zeros((nb_patches, nb_patches), np.float32)
        for i, patch_i in enumerate(self):
            for j, patch_j in enumerate(self[i + 1 :], start=i + 1):
                inter_area = (patch_i.domain & patch_j.domain).area
                areas[i, j] = inter_area
                areas[j, i] = inter_area

        if threshold > 0:
            for i in range(nb_patches):
                for j in range(nb_patches):
                    if math.sqrt(areas[i, j]) < threshold * math.sqrt(min(areas[i].max(), areas[:, j].max())):
                        areas[i, j] = 0

        G = nx.from_numpy_array(areas)
        return G

    def patch_optimal_dependency_graph(self):
        G = self.patch_intersection_graph()
        return nx.maximum_spanning_tree(G)

    def stitch(self, blend: BlendMode = "gaussian", return_mask=False):
        """
        Stitch the patches together into a single image.

        Parameters
        ----------
        blend : bool
            Whether to blend the patches together using their pixel reliability.

        Returns
        -------
        np.ndarray
            The stitched image.

        """
        shape, patch_domains = self.crop_to_content(inplace=False)
        stitch = np.zeros(shape + (3,), dtype=np.float32)
        weight = np.zeros(shape, dtype=np.float32) if blend else None

        if blend == "min":
            stitch += 1

        for patch, domain in zip(self, patch_domains, strict=True):
            patch_data = patch.resized_image(domain)
            match blend:
                case "last":
                    stitch[domain.slice()] = patch_data
                case "max":
                    stitch[domain.slice()] = np.maximum(stitch[domain.slice()], patch_data)
                case "min":
                    stitch[domain.slice()] = np.minimum(stitch[domain.slice()], patch_data)
                    weight[domain.slice()] = 1
                case "radial":
                    G = radial_mask(domain.shape)
                    stitch[domain.slice()] += G[..., None] * patch_data
                    weight[domain.slice()] += G
                case "gaussian":
                    G = gaussian_mask(domain.shape)
                    stitch[domain.slice()] += G[..., None] * patch_data
                    weight[domain.slice()] += G
                case "mean":
                    stitch[domain.slice()] += patch_data
                    weight[domain.slice()] += 1

        if blend in ("mean", "gaussian", "radial"):
            mask = weight > 0
            for i in range(stitch.shape[2]):
                stitch[..., i][mask] /= weight[mask]
        elif blend == "min":
            mask = weight > 0
            stitch[mask] = [0, 0, 0]

        if return_mask:
            return stitch, weight > 0
        return stitch


class MultiPatchRegistration:
    def __init__(self):
        self.R = sitk.ImageRegistrationMethod()
        self.R.SetMetricAsMattesMutualInformation()
        self.R.SetOptimizerAsRegularStepGradientDescent(
            learningRate=35.0,
            minStep=1,
            numberOfIterations=120,
            gradientMagnitudeTolerance=5e-6,
        )

    def __call__(self, patches: PatchStitching, return_debug=False, verbose=False):
        return self.one2one_registration(patches, return_debug=return_debug, verbose=verbose)

    def single_registration(self, fixed_patch: Patch | PatchStitching, moving_patch: Patch):
        """
        Perform registrations patch by patch for each edge in the graph.
        The result of the registration (relative displacement and similarity metric) is stored in the edge data.
        """
        # Extract the common area between the two patches
        inter_domain = (fixed_patch.domain & moving_patch.domain).to_int()
        if isinstance(fixed_patch, Patch):
            fixed_data = fixed_patch.extract_domain(inter_domain).resized_image()
        else:
            fixed_patch = fixed_patch.extract_domain(inter_domain)
            fixed_data = fixed_patch.stitch(blend="mean")
            inter_domain = fixed_patch.domain.to_int()
        moving_data = moving_patch.extract_domain(inter_domain).resized_image()

        # Convert the images to sitk
        sitk_fixed = sitk.GetImageFromArray(fixed_data[..., 1])  # We only use the green channel
        sitk_moving = sitk.GetImageFromArray(moving_data[..., 1])

        # Register the patch to the stitched image
        transform = sitk.TranslationTransform(2, [0, 0])
        self.R.SetInitialTransform(transform)

        transform = self.R.Execute(sitk_fixed, sitk_moving)
        dx, dy = transform.GetOffset()
        metric = self.R.GetMetricValue()
        return Point(-dy, -dx), metric

    def one2one_registration(self, patches: PatchStitching, area_threshold=0.35, return_debug=False, verbose=False):
        """
        Register multiple patches together.
            - First, each pair of patches with a sufficient overlap are registered together in order to construct a
              graph of relative offset between patches. The similarity metric of each registration is also stored in
              the edge data.
            - When many patches overlap together, only a subset of the relative offsets must be kept. The optimal
              dependency graph is computed as the spanning tree which maximise the total similarity metric.
            - The patches are then registered together in the order given by the optimal dependency graph.

        """
        debug = {}

        # Compute the order in which the patches should be registered
        reg_G = patches.patch_intersection_graph(area_threshold)

        # Perform the registration for each sub-graph of patches
        for p1, p2 in tqdm(reg_G.edges, desc="Registering patches two by two...", disable=not verbose, leave=False):
            try:
                dpos, metric = self.single_registration(patches[p1], patches[p2])
            except RuntimeError:
                reg_G.edges[p1, p2]["metric"] = 0
                reg_G.edges[p1, p2]["dyx"] = Point(0, 0)
            else:
                reg_G.edges[p1, p2]["dyx"] = dpos
                reg_G.edges[p1, p2]["metric"] = -metric
            reg_G.edges[p1, p2]["moving"] = p2

        debug["registration_graph"] = reg_G

        # Find the optimal reconstruction path
        optimal_path = nx.maximum_spanning_tree(reg_G, weight="metric")
        optimal_metric = sum(reg_G.edges[e]["metric"] for e in optimal_path.edges)
        if verbose:
            print(f"\tTotal metric of optimal path: {optimal_metric:.2f}")
        debug["optimal_path"] = optimal_path
        debug["optimal_metric"] = optimal_metric

        # Reconstruct the stitched image
        registered_patches = PatchStitching()
        relative_offset = {}

        for subgraph_ids in nx.connected_components(optimal_path):
            first_node = next(iter(subgraph_ids))
            if len(subgraph_ids) == 1:
                registered_patches.append(patches[first_node])
                continue

            path = optimal_path.subgraph(subgraph_ids)

            # Compute the relative displacement of each patch
            patches_dpos = {first_node: Point(0, 0)}
            for p1, p2 in nx.bfs_edges(path, source=first_node):
                dpos = reg_G.edges[p1, p2]["dyx"]
                if reg_G.edges[p1, p2]["moving"] == p1:
                    dpos = -dpos
                patches_dpos[p2] = patches_dpos[p1] + dpos

            # Center the displacements
            avg_dpos = sum(patches_dpos.values()) / len(patches_dpos)
            avg_dpos = avg_dpos.to_int()

            # Apply the displacements
            for patch_id, dpos in patches_dpos.items():
                patch = patches[patch_id].duplicate()
                offset = dpos - avg_dpos
                patch.domain = patch.domain + offset
                registered_patches.append(patch)
                relative_offset[patch_id] = offset / patch.domain.shape

        debug["relative_offset"] = relative_offset

        if return_debug:
            return registered_patches, debug
        return registered_patches

    def one2all_registration(self, patches: PatchStitching, assemble_registered_patches=True):
        # Compute the order in which the patches should be registered
        reg_G = patches.patch_intersection_graph()
        G = nx.maximum_spanning_tree(reg_G)

        # Perform the registration for each sub-graph of patches
        global_stitch = PatchStitching([])

        for g in nx.connected_components(G):
            patch0 = next(iter(g))
            if len(g) == 1:
                global_stitch.append(patches[patch0])
                continue

            g = G.subgraph(g)
            # Find the root of the tree
            patch0 = sorted(nx.center(g), key=lambda i: g.degree(i), reverse=True)[0]
            # Find the order in which the patches should be registered
            patch_seq = [edge[1] for ]

            partial_stitch = PatchStitching([patches[patch0]])

            for p1, p2 in nx.bfs_edges(g, patch0):
                if assemble_registered_patches:
                    patch_fix = partial_stitch
                else:
                    patch_fix = patches[p1]
                patch = patches[p2]
                (dy, dx), metric = self.single_registration(patch_fix, patch)

                patch.domain = patch.domain.translate(dy, dx)
                partial_stitch.append(patch)

            global_stitch += partial_stitch
        return global_stitch


def radial_mask(shape: Tuple[int, int]) -> np.ndarray:
    y, x = np.meshgrid(np.linspace(-1, 1, shape[1]), np.linspace(-1, 1, shape[0]))
    d = np.sqrt(x * x + y * y) / np.sqrt(2)
    return np.clip(1 - d, 0, 1)


def gaussian_mask(shape: Tuple[int, int], sigma=0.3) -> np.ndarray:
    y, x = np.meshgrid(np.linspace(-1, 1, shape[1]), np.linspace(-1, 1, shape[0]))
    d = np.sqrt(x * x + y * y)
    return np.exp(-(d**2 / (2.0 * sigma**2)))
