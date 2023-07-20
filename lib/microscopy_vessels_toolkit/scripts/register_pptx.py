from collections.abc import Container  # noqa
from datetime import datetime
from os.path import splitext

import cv2
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture

from microscopy_vessels_toolkit.preprocess.stitch import MultiPatchRegistration, PatchStitching


def register_pptx(pptx_path: str, output_path: str = None):
    register = MultiPatchRegistration()
    prs = Presentation(pptx_path)

    if output_path is None:
        pptx_name, pptx_ext = splitext(pptx_path)
        output_path = pptx_name + "_registered" + pptx_ext

    output_is_ppt = splitext(output_path)[1] in (".ppt", ".pptx")

    for i, slide in enumerate(prs.slides):
        stitch = PatchStitching.from_pptx_slide(slide)
        if len(stitch) < 2:
            print(f"> Slide {i+1}: No patches found.")
            continue
        else:
            print(f"> Slide {i+1}: {len(stitch)} patches found.")

        t = datetime.now()
        stitch.crop_to_content()
        reg_stitch, reg_debug = register(stitch, return_debug=True, verbose=True)
        print(f"\t... registered in {(datetime.now() - t).total_seconds():.1f} seconds.")

        if output_is_ppt:
            relative_offset = {stitch[k].alias: v for k, v in reg_debug["relative_offset"].items()}
            update_pictures_recursive(slide.shapes, relative_offset)
        else:
            filename, ext = splitext(output_path)
            slide_output_path = filename + f"_{i+1:02d}." + ext
            print(f"\t saving slide to {slide_output_path}...")
            t = datetime.now()
            stitched_img = stitch.stitch()
            cv2.imwrite(slide_output_path, stitched_img)
            print(f"\t... saved in {(datetime.now() - t).total_seconds():.1f} seconds.")

    if output_is_ppt:
        print(f"> Saving output to {output_path}...")
        t = datetime.now()
        prs.save(output_path)
        print(f"\t... saved in {(datetime.now() - t).total_seconds():.1f} seconds.")


def update_pictures_recursive(group, relative_offset):
    for s in group:
        if isinstance(s, GroupShape):
            if s.shapes:
                # Find pictures in group
                update_pictures_recursive(s.shapes, relative_offset)
        elif isinstance(s, Picture) and s.name in relative_offset:
            s.left += int(relative_offset[s.name].x * s.width)
            s.top += int(relative_offset[s.name].y * s.height)
